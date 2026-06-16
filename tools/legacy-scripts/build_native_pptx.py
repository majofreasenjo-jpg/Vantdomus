import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.dml.color import RGBColor

prs = Presentation()
# Usar layout en blanco (Slide layout 6 suele ser blanco sin nada)
slide = prs.slides.add_slide(prs.slide_layouts[6])

# Colores (Light Theme)
solomon_color = RGBColor(16, 185, 129) # Emerald 500
pillar_color = RGBColor(59, 130, 246) # Blue 500
metric_color = RGBColor(71, 85, 105) # Slate 600
line_color = RGBColor(203, 213, 225) # Slate 300

# 1. Definir Nodos (Solomon)
# X = 0.5, y empieza en 1.5, height = 0.6, width = 2.5, spacing = 0.2
solomon_nodes = [
    {"id": "s_tup", "text": "TUP (Capacidad)"},
    {"id": "s_do", "text": "Disp. Operativa"},
    {"id": "s_iie", "text": "Intensidad Energ. (IIE)"},
    {"id": "s_opex", "text": "OPEX Base"},
    {"id": "s_man", "text": "Cx Mantenimiento"},
    {"id": "s_pro", "text": "Productividad Laboral"}
]

solomon_shapes = {}
y_offset = 1.0
t_x = Inches(0.5)
t_w = Inches(2.2)
t_h = Inches(0.7)
for n in solomon_nodes:
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, t_x, Inches(y_offset), t_w, t_h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = solomon_color
    shape.line.color.rgb = solomon_color
    text_frame = shape.text_frame
    text_frame.text = n["text"]
    text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
    text_frame.paragraphs[0].font.bold = True
    solomon_shapes[n["id"]] = shape
    y_offset += 1.0

# 2. Definir Nodos (Pillars)
pillars_nodes = [
    {"id": "p_uti", "text": "1. Utilidad", "targets": ["s_tup","s_do","s_iie","s_opex","s_man","s_pro"]},
    {"id": "p_sal", "text": "2. Vida y Salud", "targets": ["s_pro","s_opex"]},
    {"id": "p_amb", "text": "3. Medioambiental", "targets": ["s_tup","s_do","s_iie","s_opex","s_man","s_pro"]},
    {"id": "p_com", "text": "4. Comunidades", "targets": ["s_tup","s_do","s_opex","s_pro"]},
    {"id": "p_per", "text": "5. Personal", "targets": ["s_pro","s_opex"]},
    {"id": "p_aud", "text": "6. Auditorías", "targets": ["s_tup","s_do","s_iie","s_opex","s_man","s_pro"]}
]

pillar_shapes = {}
y_offset = 1.0
t_x = Inches(4.0)
t_w = Inches(2.2)
for n in pillars_nodes:
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, t_x, Inches(y_offset), t_w, t_h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = pillar_color
    shape.line.color.rgb = pillar_color
    text_frame = shape.text_frame
    text_frame.text = n["text"]
    text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
    text_frame.paragraphs[0].font.bold = True
    pillar_shapes[n["id"]] = shape
    y_offset += 1.0

# 3. Definir Nodos (Metrics)
metrics_nodes = [
    {"id": "m_u1", "text": "Costo Prod.", "y": 1.0, "p": "p_uti"},
    {"id": "m_u2", "text": "Producción", "y": 1.3, "p": "p_uti"},
    {"id": "m_u3", "text": "Fallas Ops.", "y": 1.6, "p": "p_uti"},
    {"id": "m_s1", "text": "IF / IS", "y": 2.2, "p": "p_sal"},
    {"id": "m_s2", "text": "Riesgos Altos", "y": 2.6, "p": "p_sal"},
    {"id": "m_a1", "text": "Emisiones", "y": 3.2, "p": "p_amb"},
    {"id": "m_a2", "text": "Cump. Crítico", "y": 3.6, "p": "p_amb"},
    {"id": "m_c1", "text": "Reclamos", "y": 4.2, "p": "p_com"},
    {"id": "m_c2", "text": "Cump. PAG", "y": 4.6, "p": "p_com"},
    {"id": "m_p1", "text": "Capacitación", "y": 5.2, "p": "p_per"},
    {"id": "m_p2", "text": "Sobretiempo", "y": 5.6, "p": "p_per"},
    {"id": "m_ad1", "text": "Aud. Int/MPD", "y": 6.3, "p": "p_aud"}
]

metric_shapes = {}
t_x = Inches(7.5)
t_w = Inches(2.0)
t_hm = Inches(0.25)
for n in metrics_nodes:
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, t_x, Inches(n["y"]), t_w, t_hm)
    shape.fill.solid()
    shape.fill.fore_color.rgb = metric_color
    shape.line.color.rgb = metric_color
    text_frame = shape.text_frame
    text_frame.text = n["text"]
    text_frame.paragraphs[0].font.size = Pt(10)
    text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
    metric_shapes[n["id"]] = shape

# Function to connect shapes visually (since python-pptx doesn't auto-route visually on render always)
def add_connector(shape1, shape2):
    connector = slide.shapes.add_connector(MSO_CONNECTOR.CURVE, 
                    shape1.left + shape1.width, shape1.top + shape1.height/2, 
                    shape2.left, shape2.top + shape2.height/2)
    connector.line.color.rgb = line_color
    connector.line.width = Pt(1)

# Conectar Solomon <- Pillar (Line starts at Solomon right, ends at Pillar left)
for p in pillars_nodes:
    p_shape = pillar_shapes[p["id"]]
    for s_id in p["targets"]:
        s_shape = solomon_shapes[s_id]
        # Invertimos: de Solomon a Pillar (Solomon izquierdo)
        add_connector(s_shape, p_shape)

# Conectar Pillar <- Metrics
for m in metrics_nodes:
    m_shape = metric_shapes[m["id"]]
    p_shape = pillar_shapes[m["p"]]
    # De Pillar a Metric
    add_connector(p_shape, m_shape)

# Cabeceras (Títulos de Columnas)
def add_title(text, x, y):
    txBox = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(2.5), Inches(0.5))
    tf = txBox.text_frame
    p = tf.add_paragraph()
    p.text = text
    p.font.bold = True
    p.font.size = Pt(14)
    p.font.color.rgb = RGBColor(30, 41, 59)

add_title("Métricas Solomon L1", 0.5, 0.3)
add_title("Pilares Intermedios L2", 4.0, 0.3)
add_title("Métricas Terreno L3", 7.5, 0.3)

out_file = r"C:\Users\casa\Downloads\Matriz_VantDomus_PowerPoint_Nativo.pptx"
prs.save(out_file)
print(out_file)
