import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

base_dir = r"D:\Aplicaciones de Juegos\VantDomus\Antucoya\Contrato Puma"
out_file = os.path.join(base_dir, "Presentacion_Antucoya_Ejecutiva_v2.pptx")

photos = sorted([f for f in os.listdir(base_dir) if f.startswith("WhatsApp Image")])

def add_premium_infographic_slide(prs, title, sub, bullets, img_path=None, is_critical=False):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Base background (Dark Platinum/Grey)
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = RGBColor(27, 34, 44) if not is_critical else RGBColor(43, 22, 27)
    bg.line.color.rgb = bg.fill.fore_color.rgb

    # Left Accent Sidebar
    sidebar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(0.5), prs.slide_height)
    sidebar.fill.solid()
    sidebar.fill.fore_color.rgb = RGBColor(0, 168, 204) if not is_critical else RGBColor(231, 76, 60)
    sidebar.line.fill.background()

    # Premium Header Ribbon
    ribbon = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(0), Inches(12.833), Inches(1.1))
    ribbon.fill.solid()
    ribbon.fill.fore_color.rgb = RGBColor(19, 24, 32) if not is_critical else RGBColor(30, 15, 18)
    ribbon.line.fill.background()

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.15), Inches(12), Inches(0.6))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title.upper()
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.font.name = "Segoe UI"
    
    # Subtitle with accent color
    sub_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.2), Inches(12), Inches(0.5))
    tf_sub = sub_box.text_frame
    p_sub = tf_sub.paragraphs[0]
    p_sub.text = sub
    p_sub.font.size = Pt(22)
    p_sub.font.bold = True
    p_sub.font.name = "Segoe UI Light"
    p_sub.font.color.rgb = RGBColor(0, 210, 255) if not is_critical else RGBColor(255, 100, 100)

    # Infographic Bullet blocks
    start_y = 2.0
    for i, bullet in enumerate(bullets):
        
        # Rounded Rectangle as background for each bullet
        block = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(0.8), 
            Inches(start_y), 
            Inches(6.8) if img_path else Inches(11.5), 
            Inches(1.0)
        )
        block.fill.solid()
        block.fill.fore_color.rgb = RGBColor(38, 48, 62) if not is_critical else RGBColor(60, 30, 38)
        block.line.fill.background()

        # Vertical line inside the block
        accent = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0.8),
            Inches(start_y),
            Inches(0.15),
            Inches(1.0)
        )
        accent.fill.solid()
        accent.fill.fore_color.rgb = RGBColor(0, 168, 204) if not is_critical else RGBColor(231, 76, 60)
        accent.line.fill.background()

        # Text inside block
        txt = slide.shapes.add_textbox(
            Inches(1.1),
            Inches(start_y + 0.1),
            Inches(6.4) if img_path else Inches(11.0),
            Inches(0.8)
        )
        tft = txt.text_frame
        tft.word_wrap = True
        pt = tft.paragraphs[0]
        pt.text = bullet
        pt.font.size = Pt(17)
        pt.font.color.rgb = RGBColor(240, 240, 245)
        pt.font.name = "Segoe UI"
        
        start_y += 1.3

    # Footer
    footer = slide.shapes.add_textbox(Inches(0.8), Inches(6.9), Inches(12), Inches(0.5))
    p_footer = footer.text_frame.paragraphs[0]
    p_footer.text = "MINERA ANTUCOYA | CONTRATO PUMA 724 | PRIVADO Y CONFIDENCIAL"
    p_footer.font.size = Pt(10)
    p_footer.font.name = "Segoe UI"
    p_footer.font.color.rgb = RGBColor(120, 130, 140)

    # Image (Premium framing)
    if img_path and os.path.exists(img_path):
        left = Inches(8.0)
        top = Inches(2.0)
        width = Inches(4.5)
        
        # Outer Golden/Silver Frame
        frame = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left - Inches(0.08), top - Inches(0.08), width + Inches(0.16), Inches(3.8) + Inches(0.16))
        frame.fill.solid()
        frame.fill.fore_color.rgb = RGBColor(180, 180, 190) if not is_critical else RGBColor(200, 100, 100)
        frame.line.fill.background()

        try:
            slide.shapes.add_picture(img_path, left, top, width=width)
        except Exception as e:
            print(f"Error img: {e}")


def main():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    slides_data = [
        {
            "sub": "Alcance General y Programa de Ejecución (Contrato PUMA)",
            "critical": False,
            "bullets": [
                "Intervención integral de Canaletas HDPE (Pila Este/Oeste) y reposición de paquete impermeable en Piscinas PLS/ILS.",
                "Plazo del contrato: 11-Mar-26 al 04-Abr-27.",
                "Objetivo: Garantizar estanqueidad y conductividad de fluidos minimizando el impacto en la operación de Minera Antucoya."
            ]
        },
        {
            "sub": "⚠️ RUTA CRÍTICA Y ENTREGABLES ESTRATÉGICOS",
            "critical": True,
            "bullets": [
                "El programa está fuertemente condicionado por la entrega progresiva de áreas (F1 a F4) por parte del mandante.",
                "Hito Crítico 1 (PASCINAS): Término de Reparación Piscinas PLS/ILS fijado para el 03-Feb-2027.",
                "Hito Crítico 2 (CANALETAS): Término Entubado Canaletas Lixiviación fijado para el 09-Mar-2027."
            ]
        },
        {
            "sub": "Estrategia Logística - 'Tren de Actividades'",
            "critical": False,
            "bullets": [
                "Ejecución de canaletas mediante Frentes Simultáneos desfasados (2 a 3 días).",
                "Mientras el Frente 1 estabiliza cañerías, el Frente 2 entra inmediatamente en fundaciones y compactación.",
                "Impacto Directo: Reducción estimada del plazo total por zona en un 30% a 40%, optimizando grúas y personal."
            ]
        },
        {
            "sub": "Metodología - Canaletas (Trabajos Previos y Sello)",
            "critical": False,
            "bullets": [
                "Habilitación provisoria de instalaciones logísticas y apertura de accesos rápidos para maquinaria.",
                "Retiro expedito de borra remanente y mejoramiento estructural del fondo de fundación.",
                "Estricto control topográfico de compactación de sello en la zanja para evitar hundimientos o deformaciones."
            ]
        },
        {
            "sub": "Metodología - Canaletas (Montaje HDPE Corrugado 1.5m)",
            "critical": False,
            "bullets": [
                "Posicionamiento milimétrico de cañerías utilizando camión pluma certificado.",
                "Proceso de termofusión y alineamiento preciso según las pendientes críticas de drenaje hidráulico.",
                "Restricción técnica absoluta: limpieza industrial de las áreas de contacto previo a asentar las medias cañas."
            ]
        },
        {
            "sub": "Metodología - Canaletas (Estabilización y Protección)",
            "critical": False,
            "bullets": [
                "Instalación ágil de maxisacos laterales para proveer estabilización temporal de los tubos soldados.",
                "Descarga y ejecución de camellones defensivos perimetrales, protegiendo el HDPE de perforaciones.",
                "Métrica Proyectada: Velocidad de destrabe de 10 a 11 días efectivos por cada frente de trabajo en desarrollo."
            ]
        },
        {
            "sub": "Operativa en Piscinas PLS / ILS (Desmonte Táctico)",
            "critical": False,
            "bullets": [
                "Escuadrón asignado: 14 especialistas (soldadores calificados HDPE y ayudantes de maniobras).",
                "Secuencia quirúrgica de retiro multicapa: Extracción de Geomembrana de desgaste, Geonet y Geotextil inferior.",
                "Presupuesto de Tiempo: Desanclaje, corte, dimensionamiento y retiro confinado a solo 20 días la fase completa."
            ]
        },
        {
            "sub": "Operativa en Piscinas PLS / ILS (Instalación de Revestimiento)",
            "critical": False,
            "bullets": [
                "Rendimientos de instalación en Piscinas: Geomembrana LLDPE a 2.000 m²/día, y Geosintéticos a 2.857 m²/día.",
                "Ingeniería Multicapa: Sellado mediante Geotextil (300g/m²), LLDPE, Geonet (5mm), y capa blindada HDPE (2mm).",
                "Duración proyectada: 34 días críticos por cada piscina para sello técnico (Total ciclo ~54 días; 20.000 m2)."
            ]
        },
        {
            "sub": "Aseguramiento y Control de Calidad Permanente (QA/QC)",
            "critical": False,
            "bullets": [
                "Tecnología en Terreno: Despliegue de 4 extrusoras simultáneas y 4 cuñas de soldadura térmica de autopropulsión.",
                "Inspección de Fallas Cero: Ensayos de vacío (Vacuum Box) y monitor de integridad de chispa (Spark Tester).",
                "Liberación de carpeta de Calidad garantizando validación topográfica e impermeabilidad de la red para Antucoya."
            ]
        },
        {
            "sub": "Compromiso Sustentable HSEC (Cero Daño)",
            "critical": False,
            "bullets": [
                "Condición de Operación Extrema: Uso intransable de líneas de vida en taludes y salvavidas en bordes ILS.",
                "Estrategia Anti-Fatalidades: Separación táctica Hombre-Máquina dentro de trincheras y piscinas (Rigger Activo).",
                "Alineación 100% incondicional con las Directrices de Riesgo Crítico, Fatiga y Supervisión de Minera Antucoya."
            ]
        }
    ]

    for i, slide_info in enumerate(slides_data):
        img = os.path.join(base_dir, photos[i]) if i < len(photos) else None
        add_premium_infographic_slide(
            prs,
            title="Introducción del programa de construcción",
            sub=slide_info["sub"],
            bullets=slide_info["bullets"],
            img_path=img,
            is_critical=slide_info["critical"]
        )

    prs.save(out_file)
    print(f"Saved: {out_file}")

if __name__ == "__main__":
    main()
