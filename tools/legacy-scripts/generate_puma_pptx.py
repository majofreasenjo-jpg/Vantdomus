import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

base_dir = r"D:\Aplicaciones de Juegos\VantDomus\Antucoya\Contrato Puma"
out_file = os.path.join(base_dir, "Presentacion_Antucoya_Ejecutiva.pptx")

# 15 WhatsApp images
photos = sorted([f for f in os.listdir(base_dir) if f.startswith("WhatsApp Image")])

def add_premium_slide(prs, title, sub, bullets, img_path=None, is_critical=False):
    # Layout 6 is usually BLANK
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Background Shape
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    # Dark Premium Background (Deep Navy Blue or Maroon for Critical)
    bg.fill.solid()
    if is_critical:
        bg.fill.fore_color.rgb = RGBColor(60, 15, 20)  # Dark Red
    else:
        bg.fill.fore_color.rgb = RGBColor(20, 25, 35)  # Dark Blueish Grey
    bg.line.color.rgb = bg.fill.fore_color.rgb

    # Title Box
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(12.33), Inches(1.0))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title.upper()
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.font.name = "Arial"
    
    # Subtitle / Accent Line
    # A golden/cyan thin line under title
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1.2), Inches(12.33), Inches(0.04))
    line.fill.solid()
    if is_critical:
        line.fill.fore_color.rgb = RGBColor(255, 80, 80)
    else:
        line.fill.fore_color.rgb = RGBColor(0, 200, 255) # Cyan
    line.line.fill.background()

    sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.4), Inches(12.33), Inches(0.8))
    tf_sub = sub_box.text_frame
    p_sub = tf_sub.paragraphs[0]
    p_sub.text = sub
    p_sub.font.size = Pt(20)
    p_sub.font.bold = True
    p_sub.font.color.rgb = RGBColor(200, 200, 210)
    if is_critical:
        p_sub.font.color.rgb = RGBColor(255, 200, 200)

    # Bullets Box
    # If we have an image, bullets take less width
    bullet_width = Inches(12.33) if not img_path else Inches(7.0)
    content_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.3), bullet_width, Inches(4.5))
    tf_content = content_box.text_frame
    tf_content.word_wrap = True
    
    for i, bullet in enumerate(bullets):
        p_bullet = tf_content.add_paragraph() if i > 0 else tf_content.paragraphs[0]
        p_bullet.text = "• " + bullet
        p_bullet.font.size = Pt(17)
        p_bullet.font.color.rgb = RGBColor(230, 230, 240)
        p_bullet.space_after = Pt(20)
        p_bullet.space_before = Pt(10)
        if is_critical:
            p_bullet.font.bold = True
            p_bullet.font.size = Pt(18)

    # Footer
    footer = slide.shapes.add_textbox(Inches(0.5), Inches(6.9), Inches(12.33), Inches(0.5))
    p_footer = footer.text_frame.paragraphs[0]
    p_footer.text = "MINERA ANTUCOYA | CONTRATO PUMA 724 | PRIVADO Y CONFIDENCIAL"
    p_footer.font.size = Pt(10)
    p_footer.font.color.rgb = RGBColor(100, 110, 120)

    # Add Image
    if img_path and os.path.exists(img_path):
        # We enforce strict alignment and bounding box
        left = Inches(7.8)
        top = Inches(2.3)
        width = Inches(4.8)
        
        # We add a subtle glow/border behind the image
        border = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left - Inches(0.05), top - Inches(0.05), width + Inches(0.1), Inches(3.6) + Inches(0.1))
        border.fill.solid()
        border.fill.fore_color.rgb = RGBColor(255, 255, 255) if not is_critical else RGBColor(255, 100, 100)
        border.line.fill.background()
        
        try:
            # We fix the width, PowerPoint will scale height to keep ratio
            # This ensures they are perfectly aligned on the left/right axis
            pic = slide.shapes.add_picture(img_path, left, top, width=width)
        except Exception as e:
            print(f"Error {e}")


def main():
    prs = Presentation()
    # WIDESCREEN 16:9 ASPECT RATIO
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    slides_data = [
        {
            "sub": "Alcance General y Programa de Ejecución (Contrato PUMA)",
            "critical": False,
            "bullets": [
                "Intervención integral de Canaletas HDPE (Pila Este/Oeste) y reposición de paquete impermeable en Piscinas PLS/ILS.",
                "Plazo del contrato: 11-Mar-26 al 04-Abr-27.",
                "Objetivo: Garantizar estanqueidad y conductividad de fluidos minimizando el impacto en la operación de Minera Antucoya."
            ]
        },
        {
            "sub": "⚠️ RUTA CRÍTICA Y ENTREGABLES ESTRATÉGICOS",
            "critical": True,
            "bullets": [
                "El programa está fuertemente condicionado por la entrega progresiva de áreas (F1 a F4) por parte del mandante.",
                "Hito 1: Entrega de Área F1 Pila Dinámica Este: 14-Abr-26.",
                "Hito Crítico 2 (PISICNAS): Término de Reparación Piscinas PLS/ILS: 03-Feb-2027.",
                "Hito Crítico 3 (CANALETAS): Término Entubado Canaletas Lixiviación: 09-Mar-2027.",
                "El cumplimiento exige la movilización acelerada antes del 14-Abr-26 para mitigar cuellos de botella."
            ]
        },
        {
            "sub": "Estrategia Logística - 'Tren de Actividades'",
            "critical": False,
            "bullets": [
                "Ejecución de canaletas mediante Frentes Simultáneos desfasados (2 a 3 días).",
                "Mientras el Frente 1 estabiliza cañerías, el Frente 2 entra inmediatamente en fundaciones y compactación.",
                "Impacto Directo: Reducción estimada del plazo total por zona en un 30% a 40%, optimizando grúas y personal."
            ]
        },
        {
            "sub": "Metodología - Canaletas (Trabajos Previos y Sello)",
            "critical": False,
            "bullets": [
                "Habilitación provisoria de instalaciones logísticas y apertura de accesos rápidos para maquinaria.",
                "Retiro expedito de borra remanente y mejoramiento de fondo de fundación.",
                "Estricto aseguramiento y control topográfico de compactación de sello en la zanja para evitar hundimientos o deformaciones en la matriz de HDPE."
            ]
        },
        {
            "sub": "Metodología - Canaletas (Montaje HDPE Corrugado 1.5m)",
            "critical": False,
            "bullets": [
                "Posicionamiento milimétrico de cañerías utilizando camión pluma certificado.",
                "Proceso térmico de acople (fusión o unión mecánica) con cuadrillas especializadas y alineación precisa de pendientes de drenaje.",
                "Restricción técnica absoluta: limpieza industrial de las áreas de contacto previo a asentar las medias cañas."
            ]
        },
        {
            "sub": "Metodología - Canaletas (Estabilización y Protección)",
            "critical": False,
            "bullets": [
                "Instalación ágil de maxisacos laterales para proveer estabilización temporal de los tubos soldados.",
                "Descarga y ejecución de camellones defensivos con arena fina y dócil que protege el HDPE de perforaciones.",
                "Métrica Proyectada: Velocidad de destrabe de 10 a 11 días efectivos por cada frente de trabajo en desarrollo lineal."
            ]
        },
        {
            "sub": "Operativa en Piscinas PLS / ILS (Desmonte Táctico)",
            "critical": False,
            "bullets": [
                "Escuadrón asignado: 14 especialistas (soldadores calificados HDPE y ayudantes de maniobras).",
                "Secuencia quirúrgica de retiro multicapa: Extracción de Geomembrana (Etapa 1), seguida de Geonet estructural y Geotextil inferior.",
                "Presupuesto de Tiempo: Desanclaje, corte, dimensionamiento y retiro confinado a solo 5 días por capa (Total 20 días la fase completa)."
            ]
        },
        {
            "sub": "Operativa en Piscinas PLS / ILS (Instalación de Revestimiento Multicapa)",
            "critical": False,
            "bullets": [
                "Rendimiento agresivo LLDPE programado: Avance de ~2.000 m² hasta picos de 2.857 m² instalados por día.",
                "Ingeniería Multicapa Inversa: Sellado mediante Geotextil de 300g/m², Geomembrana base LLDPE, Red Drenaje Geonet de 5mm, y capa blindada HDPE 2mm exterior.",
                "Duración proyectada: 34 días críticos por cada piscina logrando estanqueidad hermética (Total del ciclo ~54 días con revestimiento 20.000m2)."
            ]
        },
        {
            "sub": "Aseguramiento y Control de Calidad Transversal (QA/QC)",
            "critical": False,
            "bullets": [
                "Tecnología en Terreno: Despliegue de 4 extrusoras simultáneas, 4 cuñas de soldadura térmica de autopropulsión, y manómetros de aguja de precisión.",
                "Inspección de Fallas Cero: Pruebas mecánicas diarias con tensiómetros, ensayos de vacío (Vacuum Box) y monitor de integridad de chispa (Spark Tester) constante.",
                "Dossier Técnico: Liberación de carpeta de Calidad garantizando la impermeabilidad de la red PLS/ILS para Minera Antucoya."
            ]
        },
        {
            "sub": "Compromiso HSEC (Estándar Cero Daño Permanente)",
            "critical": False,
            "bullets": [
                "Condición de Operación Extrema: Uso intransable de líneas de vida en taludes, salvavidas en bordes ILS, y PPE de Categoría Superior contra químicos y cortes.",
                "Estrategia Anti-Fatalidades: Segmentación táctica minimizando la exposición de obreros a la maquinaria pesada dentro de las trincheras y piscinas (Rigger Activo).",
                "Integración Absoluta: Alineación 100% incondicional con las Directrices de Riesgo Crítico, Fatiga y Supervisión Estratégica mandatadas por el ecosistema de Antucoya."
            ]
        }
    ]

    for i, slide_info in enumerate(slides_data):
        img_path = os.path.join(base_dir, photos[i]) if i < len(photos) else None
        add_premium_slide(
            prs=prs,
            title="Introducción del programa de construcción",
            sub=slide_info["sub"],
            bullets=slide_info["bullets"],
            img_path=img_path,
            is_critical=slide_info["critical"]
        )

    prs.save(out_file)
    print(f"Premium Presentation saved to {out_file}")

if __name__ == "__main__":
    main()
