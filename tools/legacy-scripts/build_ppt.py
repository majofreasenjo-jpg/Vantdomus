from pptx import Presentation
from pptx.util import Inches, Pt

prs = Presentation()

# Slide 1: Title
title_slide_layout = prs.slide_layouts[0]
slide = prs.slides.add_slide(title_slide_layout)
title = slide.shapes.title
subtitle = slide.placeholders[1]
title.text = "Estructura Estratégica: KPIs 2026"
subtitle.text = "Análisis Jerárquico Solomon, Generado por VantDomus"

# Slide 2: Drivers Económicos y Operativos
bullet_slide_layout = prs.slide_layouts[1]
slide = prs.slides.add_slide(bullet_slide_layout)
shapes = slide.shapes
title_shape = shapes.title
body_shape = shapes.placeholders[1]
title_shape.text = "Drivers Económicos y Operativos"
tf = body_shape.text_frame
# Pilar Utilidad
p = tf.paragraphs[0]
p.text = "Pilar: Utilidad"
p.font.bold = True
p = tf.add_paragraph()
p.text = "Métricas Base: Costo Producción, Volúmenes y Detenciones (Fallas de Ops/Equipos)."
p.level = 1
p = tf.add_paragraph()
p.text = "Impacto Sistémico: Alimenta integralmente al TUP, Disponibilidad Métrica, EII y Costos Base."
p.level = 1
# Pilar Personal
p = tf.add_paragraph()
p.text = "Pilar: Personal"
p.font.bold = True
p.level = 0
p = tf.add_paragraph()
p.text = "Métricas Base: Cobertura de Entrenamiento y Control de Sobretiempo."
p.level = 1
p = tf.add_paragraph()
p.text = "Impacto Sistémico: Conductor primario de la Productividad Laboral y modulador de Margen OPEX."
p.level = 1

# Slide 3: Sostenibilidad
slide = prs.slides.add_slide(bullet_slide_layout)
shapes = slide.shapes
title_shape = shapes.title
body_shape = shapes.placeholders[1]
title_shape.text = "Sostenibilidad, Entorno y Confiabilidad (HSE)"
tf = body_shape.text_frame
# Cuidado Vida
p = tf.paragraphs[0]
p.text = "Pilar: Cuidado Vida y Salud"
p.font.bold = True
p = tf.add_paragraph()
p.text = "Métricas Base: Frecuencia/Severidad (IF/IS) y Prevención de Riesgos de Fatalidad."
p.level = 1
p = tf.add_paragraph()
p.text = "Impacto Sistémico: Golpea financieramente en el OPEX (Non-Energy) y paraliza Productividad Laboral."
p.level = 1
# Cuidado Ambiental
p = tf.add_paragraph()
p.text = "Pilar: Cuidado Ambiental & Comunidades"
p.font.bold = True
p.level = 0
p = tf.add_paragraph()
p.text = "Métricas Base: Emisiones Físicas (SO2, MP, NOX) y Tasa de Conflictividad (Reclamos)."
p.level = 1
p = tf.add_paragraph()
p.text = "Impacto Sistémico: Su incumplimiento genera restricciones que limitan asimétricamente el TUP y la Disponibilidad."
p.level = 1

# Slide 4: Governance
slide = prs.slides.add_slide(bullet_slide_layout)
shapes = slide.shapes
title_shape = shapes.title
body_shape = shapes.placeholders[1]
title_shape.text = "Governance y Control Estructural"
tf = body_shape.text_frame
p = tf.paragraphs[0]
p.text = "Pilar: Auditorías"
p.font.bold = True
p = tf.add_paragraph()
p.text = "Métricas Base: Cumplimiento de Auditorías Internas y Directrices MPD."
p.level = 1
p = tf.add_paragraph()
p.text = "Impacto Sistémico: Posee propagación 360°, afectando transversalmente la totalidad de referenciales Solomon (EII, TUP, DO y OPEX)."
p.level = 1

# Slide 5: Diagram Code
slide = prs.slides.add_slide(bullet_slide_layout)
shapes = slide.shapes
title_shape = shapes.title
body_shape = shapes.placeholders[1]
title_shape.text = "Código Mermaid del Diagrama Analítico"
tf = body_shape.text_frame
p = tf.paragraphs[0]
p.text = "Renderice este código en www.mermaid.live o pegue el PNG directo en su PPT Final:"
p = tf.add_paragraph()
p.text = "graph LR;"
p.font.size = Pt(12)
p = tf.add_paragraph()
p.text = "  Nivel3[Métricas Base] --> Nivel2[Pilares Refinería]"
p.font.size = Pt(12)
p = tf.add_paragraph()
p.text = "  Nivel2 --> Nivel1[KPIs Solomon]"
p.font.size = Pt(12)

prs.save(r"C:\Users\casa\Downloads\Estructura_KPIS_2026_VantDomus.pptx")
print("PPT Guardado")
