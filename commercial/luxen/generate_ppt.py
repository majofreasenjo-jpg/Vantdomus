import os
import subprocess
import sys

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
except ImportError:
    subprocess.check_call([sys.executable, "-m", "pip", "install", "python-pptx"])
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN

def create_vantdomus_ppt(output_path):
    prs = Presentation()
    
    # -----------------------------
    # VANTDOMUS THEME COLORS
    # -----------------------------
    DARK_NAVY = RGBColor(11, 20, 38)     # Fondo oscuro corporativo
    BRIGHT_AMBER = RGBColor(245, 158, 11) # Acentos (Bronze/Amber)
    SLATE_LIGHT = RGBColor(203, 213, 225) # Texto secundario
    WHITE = RGBColor(255, 255, 255)      # Texto primario
    
    # helper for dark background
    def set_slide_background(slide):
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = DARK_NAVY

    # SLIDE 1: PORTADA
    slide_layout = prs.slide_layouts[0] # Title slide
    slide = prs.slides.add_slide(slide_layout)
    set_slide_background(slide)
    
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    title.text = "VANTDOMUS B2B\nMODELO COMERCIAL TIER-1"
    title.text_frame.paragraphs[0].font.color.rgb = BRIGHT_AMBER
    title.text_frame.paragraphs[0].font.name = 'Arial'
    title.text_frame.paragraphs[0].font.bold = True
    
    subtitle.text = "Estructura de Captura de Valor y Arbitraje Algorítmico\nLuxen Analytics"
    subtitle.text_frame.paragraphs[0].font.color.rgb = SLATE_LIGHT

    # SLIDE 2: DOCTRINA COMERCIAL
    slide_layout = prs.slide_layouts[1] # Title & Content
    slide = prs.slides.add_slide(slide_layout)
    set_slide_background(slide)
    
    title, body = slide.shapes.title, slide.placeholders[1]
    title.text = "LA DOCTRINA ZERO-TRUST"
    title.text_frame.paragraphs[0].font.color.rgb = BRIGHT_AMBER
    
    tf = body.text_frame
    tf.text = "Nuestra postura comercial rechaza el modelo SaaS convencional."
    tf.paragraphs[0].font.color.rgb = WHITE
    
    p = tf.add_paragraph()
    p.text = "No vendemos 'Software por Usuario'."
    p.font.color.rgb = SLATE_LIGHT
    
    p = tf.add_paragraph()
    p.text = "Operamos como una Firma de Defensa Analítica."
    p.font.color.rgb = WHITE
    
    p = tf.add_paragraph()
    p.text = "El financiamiento se ancla matemáticamente a la magnitud del CAPEX y al dinero que salvamos en arbitrajes."
    p.font.color.rgb = SLATE_LIGHT

    # SLIDE 3: BRAZO 1
    slide = prs.slides.add_slide(slide_layout)
    set_slide_background(slide)
    title, body = slide.shapes.title, slide.placeholders[1]
    
    title.text = "BRAZO 1: RETAINER (EL ESCUDO ESTÁTICO)"
    title.text_frame.paragraphs[0].font.color.rgb = BRIGHT_AMBER
    
    tf = body.text_frame
    tf.text = "Cubrimos el costo operativo y bloqueamos el riesgo."
    tf.paragraphs[0].font.color.rgb = WHITE
    
    p = tf.add_paragraph()
    p.text = "Fee de Piloto Forense (FDE): 500 a 1.200 UF."
    p.font.color.rgb = SLATE_LIGHT
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "Fee VantDomus Anual: 0.25% del CAPEX del Proyecto."
    p.font.color.rgb = SLATE_LIGHT
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "Justificación: El cliente compra la inyección técnica y el análisis forense de la 'Materia Oscura' independientemente del final del arbitraje."
    p.font.color.rgb = WHITE
    p.font.bold = True
    
    # SLIDE 4: BRAZO 2
    slide = prs.slides.add_slide(slide_layout)
    set_slide_background(slide)
    title, body = slide.shapes.title, slide.placeholders[1]
    
    title.text = "BRAZO 2: SUCCESS FEE (LA EXTRACCIÓN)"
    title.text_frame.paragraphs[0].font.color.rgb = BRIGHT_AMBER
    
    tf = body.text_frame
    tf.text = "La escalabilidad real del negocio corporativo."
    tf.paragraphs[0].font.color.rgb = WHITE
    
    p = tf.add_paragraph()
    p.text = "Retención Asimétrica: 8% al 12% Neto de todo capital extraído."
    p.font.color.rgb = SLATE_LIGHT
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "Aplica para: Multas revertidas, Discrepancias de Facturación (Claims NOC), Prevención probada de sobrecostos constructivos."
    p.font.color.rgb = SLATE_LIGHT
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "Psicología: VantDomus se convierte en la inversión más lógica del Holding; el sistema se paga infinitas veces con capital recuperado."
    p.font.color.rgb = WHITE
    
    # SLIDE 5: CASO PUMA
    slide = prs.slides.add_slide(slide_layout)
    set_slide_background(slide)
    title, body = slide.shapes.title, slide.placeholders[1]
    
    title.text = "EJEMPLO: LITIGIOS MINEROS (CASO PUMA)"
    title.text_frame.paragraphs[0].font.color.rgb = BRIGHT_AMBER
    
    tf = body.text_frame
    tf.text = "Escenario: 8.603 HHs fantasmas + 54 Días Overhead (Defensa ~21.000 UF)"
    tf.paragraphs[0].font.color.rgb = WHITE
    tf.paragraphs[0].font.bold = True
    
    p = tf.add_paragraph()
    p.text = "Cobro Histórico Errados: Tarifa Plana por Horas ($2.5M CLP)."
    p.font.color.rgb = RGBColor(239, 68, 68) # Redish
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "Cobro Estructura VantDomus Titán:"
    p.font.color.rgb = WHITE
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "Retainer Setup Forense: 1.000 UF ($38 Millones CLP)."
    p.font.color.rgb = BRIGHT_AMBER
    p.level = 2
    
    p = tf.add_paragraph()
    p.text = "Plus Success Fee: 10% sobre las 21.000 UF = Extra 2.100 UF ($80 Millones CLP)."
    p.font.color.rgb = BRIGHT_AMBER
    p.level = 2
    
    # SLIDE 6: CIERRE
    slide = prs.slides.add_slide(slide_layout)
    set_slide_background(slide)
    title, body = slide.shapes.title, slide.placeholders[1]
    
    title.text = "REGLAS DE NEGOCIACIÓN B2B"
    title.text_frame.paragraphs[0].font.color.rgb = BRIGHT_AMBER
    
    tf = body.text_frame
    tf.text = "1. El precio no se baja. Se audita lo que ya perdieron esta semana y se compara."
    tf.paragraphs[0].font.color.rgb = WHITE
    
    p = tf.add_paragraph()
    p.text = "2. Somos un equipo de Extracción, no un Software contable barato."
    p.font.color.rgb = SLATE_LIGHT
    
    p = tf.add_paragraph()
    p.text = "3. Cada informe forense es Inapelable Matemáticamente, lo que justifica las penalidades."
    p.font.color.rgb = WHITE

    # Replace solid background elements inside master logic if necessary to fix unreadable text on bright elements
    for sld in prs.slides:
        for shape in sld.shapes:
            if not shape.has_text_frame: continue
            shape.text_frame.paragraphs[0].font.color.rgb = shape.text_frame.paragraphs[0].font.color.rgb or WHITE

    prs.save(output_path)
    print(f"PPT Guardado exitosamente en: {output_path}")

if __name__ == "__main__":
    out_pth = r"d:\Aplicaciones de Juegos\VantDomus\Luxen_Comercial\VantDomus_Pricing_Titan.pptx"
    create_vantdomus_ppt(out_pth)
